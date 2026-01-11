# -*- coding: utf-8 -*-
import os
import json
from typing import Any, Dict, List, Tuple, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def _emu(v: int) -> float:
    return float(v or 0)

def _clamp01(x: float) -> float:
    return max(0.0, min(1.0, x))

def _bbox_norm(left, top, width, height, sw, sh):
    x0 = _clamp01(_emu(left) / sw)
    y0 = _clamp01(_emu(top) / sh)
    x1 = _clamp01((_emu(left) + _emu(width)) / sw)
    y1 = _clamp01((_emu(top) + _emu(height)) / sh)
    return [x0, y0, x1, y1]

def _iter_shapes(shape) -> List[Any]:
    """
    Recursively yield shapes including inside groups.
    """
    out = [shape]
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            out.extend(_iter_shapes(s))
    return out

def extract_text_boxes(pptx_path: str) -> Dict[str, Any]:
    prs = Presentation(pptx_path)
    sw = float(prs.slide_width)
    sh = float(prs.slide_height)

    slides_out = []
    for si, slide in enumerate(prs.slides, start=1):
        tbs = []
        flat = []
        for shp in slide.shapes:
            flat.extend(_iter_shapes(shp))

        for zi, shp in enumerate(flat):
            # python-pptx has text_frame on many types; we only keep meaningful text
            txt = ""
            try:
                if hasattr(shp, "text_frame") and shp.text_frame:
                    txt = (shp.text_frame.text or "").strip()
            except Exception:
                txt = ""
            if not txt:
                continue

            try:
                bb = _bbox_norm(shp.left, shp.top, shp.width, shp.height, sw, sh)
            except Exception:
                bb = [0.0, 0.0, 1.0, 1.0]

            tbs.append({
                "z": zi,
                "text": txt,
                "rel_bbox": bb,
            })

        slide_texts = []
        # convenience: concat textboxes as "texts"
        for tb in sorted(tbs, key=lambda x: (x["rel_bbox"][1], x["rel_bbox"][0])):
            slide_texts.append(tb["text"])

        slides_out.append({
            "slide": si,
            "text_boxes": tbs,
            "texts": slide_texts,
        })

    return {
        "pptx": os.path.basename(pptx_path),
        "slide_width_emu": sw,
        "slide_height_emu": sh,
        "slides": slides_out,
    }

def main():
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--out", required=True)
    args = ap.parse_args()
    data = extract_text_boxes(args.pptx)
    os.makedirs(os.path.dirname(args.out), exist_ok=True)
    with open(args.out, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print("Wrote", args.out)

if __name__ == "__main__":
    main()
